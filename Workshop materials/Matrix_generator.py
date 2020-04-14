#%% Load packages and pp presentation
from pptx import Presentation
import plotly.graph_objects as go
import plotly.io as pio
pio.renderers.default = "browser"
prs = Presentation("Teknologikatalog.pptx")


#%% Get data from pp
value = []
dataReadiness = []
technology = []

def get_values(values,dataReadiness,technology,i=0):
        values.append(int(slide.shapes[i+3].text_frame.paragraphs[0].runs[0].text.replace('.','')))
        dataReadiness.append(int(slide.shapes[i+8].text_frame.paragraphs[0].runs[0].text.replace('.','')))  
        technology.append(''.join([run.text for run in slide.shapes[0].text_frame.paragraphs[0].runs]))
        return values,dataReadiness,technology

for slide in prs.slides:
    try:
        value,dataReadiness,technology = get_values(value,dataReadiness,technology)
    except :
        try : 
            value,dataReadiness,technology = get_values(value,dataReadiness,technology,1)
        except:
            pass
# %% Plot data
fig = go.Figure()
i = 1
for x,y,name in zip(dataReadiness,value,technology):
    fig.add_trace(go.Scatter(
                        x=[x],
                        y=[y],
                        mode='markers+text',
                        hovertext = name,
                        #text=[name],
                        name=name,
                        marker=dict(size=20,symbol=i%3),
                        textposition="bottom right"
                        ))
    i = i+1

fig.update_layout(
    title="Technology matrix",
    xaxis_title="Data readines",
    yaxis_title="Value",
    font=dict(
        family="Courier New, monospace",
        size=22,
        color="#7f7f7f"
    )
)
fig.show()

# %%
